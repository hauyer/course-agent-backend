import re
import unicodedata
from pathlib import Path

from charset_normalizer import from_bytes
from docx import Document
from pypdf import PdfReader
from pptx import Presentation

try:
    import fitz  # PyMuPDF，优先保留复杂 PDF 的阅读顺序和特殊字体映射
except ImportError:  # 允许旧环境先使用 pypdf 回退
    fitz = None



class DocumentParseError(Exception):
    """资料解析异常"""


def _pdf_text_quality(text: str) -> tuple[bool, str]:
    """判断页面文本层是否足够可靠，并返回面向用户的原因。"""
    compact = "".join((text or "").split())
    if len(compact) < 20:
        return False, "提取文字过少，页面可能是扫描图片"
    broken = sum(
        1
        for ch in compact
        if ch in {"\ufffd", "\u25a0", "\u25a1"} or "\ue000" <= ch <= "\uf8ff"
    )
    if broken / max(len(compact), 1) > 0.015:
        return False, "检测到较多字体映射损坏字符"
    readable = sum(ch.isalnum() or "\u4e00" <= ch <= "\u9fff" for ch in compact)
    if readable / max(len(compact), 1) < 0.35:
        return False, "可识别文字比例偏低，公式或特殊字体可能丢失"
    return True, ""


def _quality_notice(page_number: int, reason: str, ocr_attempted: bool) -> str:
    suffix = "，已尝试 OCR" if ocr_attempted else "，建议安装 Tesseract OCR 后重新解析"
    return f"[文本质量提示：第{page_number}页文本层质量较差（{reason}{suffix}）]"



def normalize_text(text: str) -> str:
    """
    对解析后的文本进行基础清洗。

    只清理多余空格和过多空行，
    尽量保留原文的段落结构。
    """

    text = unicodedata.normalize("NFKC", text or "")
    text = text.replace("\r\n","\n").replace("\r","\n")
    text = text.replace("\u00a0"," ")
    text = text.replace("\u3000"," ")

    # PDF 字体映射损坏时 pypdf 会产出替换字符、空方框或私用区字形。
    # 这些字符无法还原原始字义，但继续展示会形成大段“□□”；统一移除并
    # 在旧索引的检索出口再次执行本清洗，避免用户必须重新上传资料。
    text = re.sub(r"[\ufffd\u25a0\u25a1\ue000-\uf8ff]+", " ", text)
    text = "".join(
        ch
        for ch in text
        if ch in "\n\t" or unicodedata.category(ch) not in {"Cc", "Cf", "Cs"}
    )


    # 清理多余空格
    lines = [line.strip() for line in text.split("\n")]
    lines = [re.sub(r"[ \t]{2,}", " ", line) for line in lines]

    text = "\n".join(lines)

    # 三个及以上的空行替换为两个空行
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text.strip()


def parse_pdf(file_path:Path) ->str:
    """
    解析普通文本型 PDF。

    注意：
    扫描版 PDF 没有文本层时，pypdf 可能提取不到内容。
    后续可再增加 OCR 功能。
    """

    # PyMuPDF 对阅读顺序、嵌入字体和数学符号的处理通常优于 pypdf。
    # 文本质量不足时尝试其 OCR 接口；本机未安装 Tesseract 时给出明确提示。
    if fitz is not None:
        try:
            document = fitz.open(file_path)
            page_texts = []
            for page_number, page in enumerate(document, start=1):
                content = page.get_text("text", sort=True) or ""
                good, reason = _pdf_text_quality(content)
                ocr_attempted = False
                if not good:
                    try:
                        text_page = page.get_textpage_ocr(
                            language="chi_sim+eng", dpi=180, full=True
                        )
                        ocr_text = page.get_text(
                            "text", textpage=text_page, sort=True
                        ) or ""
                        ocr_attempted = True
                        ocr_good, ocr_reason = _pdf_text_quality(ocr_text)
                        if len(ocr_text.strip()) > len(content.strip()):
                            content = ocr_text
                        good, reason = ocr_good, ocr_reason or reason
                    except Exception:
                        pass
                content = normalize_text(content)
                blocks = [f"[第{page_number}页]"]
                if not good:
                    blocks.append(_quality_notice(page_number, reason, ocr_attempted))
                if content:
                    blocks.append(content)
                page_texts.append("\n".join(blocks))
            document.close()
            result = normalize_text("\n\n".join(page_texts))
            if result:
                return result
        except Exception:
            # 文件损坏或 PyMuPDF 不兼容时继续使用 pypdf，而不是直接使上传失败。
            pass

    try:
        reader = PdfReader(file_path)
    except Exception as exc:
        raise DocumentParseError(f"无法解析 PDF 文件: {exc}") from exc

    if reader.is_encrypted:
        try:
            decrypt_result = reader.decrypt("")
            if decrypt_result == 0:
                raise DocumentParseError("PDF 文件已加密，无法解析。")
        except Exception as exc:
            raise DocumentParseError(f"无法解密 PDF 文件: {exc}") from exc

    page_texts = []

    for page_number, page in enumerate(reader.pages, start=1):
        try:
            content = page.extract_text() or ""
        except Exception:
            content = ""

        good, reason = _pdf_text_quality(content)
        content = normalize_text(content)

        blocks = [f"[第{page_number}页]"]
        if not good:
            blocks.append(_quality_notice(page_number, reason, False))
        if content:
            blocks.append(content)
        page_texts.append("\n".join(blocks))

    return normalize_text("\n\n".join(page_texts))



def parse_docx(file_path:Path) ->str:
    """
    解析 Word DOCX 文档，包括普通段落和表格。
    """

    try:
        document = Document(str(file_path))
    except Exception as exc:
        raise DocumentParseError(f"无法解析 DOCX 文件: {exc}") from exc

    contents = []

    for paragraph in document.paragraphs:
        text = normalize_text(paragraph.text)
        if text:
            contents.append(text)

    for table_index, table in enumerate(document.tables, start=1):
        table_rows = []

        for row in table.rows:
            cells = [
                normalize_text(cell.text)
                for cell in row.cells
            ]

            if any(cells):
                table_rows.append(" | ".join(cells))

        if table_rows:
            contents.append(
                f"[表格{table_index}]\n"
                + "\n".join(table_rows)
            )

    return normalize_text("\n\n".join(contents))


def parse_pptx(file_path: Path) -> str:
    """
    解析 PowerPoint PPTX 中各页可见文本。
    """
    try:
        presentation = Presentation(str(file_path))
    except Exception as exc:
        raise DocumentParseError(f"PPTX 文件读取失败：{exc}") from exc

    slide_contents = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):
        texts = []

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue

            text = normalize_text(shape.text)

            if text:
                texts.append(text)

        if texts:
            slide_contents.append(
                f"[第{slide_number}页]\n"
                + "\n".join(texts)
            )

    return normalize_text("\n\n".join(slide_contents))


def parse_text_file(file_path: Path) -> str:
    """
    解析 TXT 或 Markdown 文件。

    使用 charset-normalizer 自动识别 UTF-8、GBK 等编码。
    """
    try:
        data = file_path.read_bytes()
    except Exception as exc:
        raise DocumentParseError(f"文本文件读取失败：{exc}") from exc

    detected = from_bytes(data).best()

    if detected is not None:
        text = str(detected)
        return normalize_text(text)

    # 自动检测失败时进行常见编码回退
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "gbk"):
        try:
            return normalize_text(data.decode(encoding))
        except UnicodeDecodeError:
            continue

    raise DocumentParseError("无法识别文本文件编码")


def parse_document(file_path: Path) -> str:
    """
    根据文件扩展名调用对应的解析器。
    """
    if not file_path.exists():
        raise DocumentParseError("资料文件不存在")

    extension = file_path.suffix.lower()

    parser_map = {
        ".pdf": parse_pdf,
        ".docx": parse_docx,
        ".pptx": parse_pptx,
        ".txt": parse_text_file,
        ".md": parse_text_file,
    }

    parser = parser_map.get(extension)

    if parser is None:
        raise DocumentParseError(
            f"暂不支持解析此文件类型：{extension}"
        )

    text = parser(file_path)

    if not text.strip():
        raise DocumentParseError(
            "未从文件中提取到有效文本。"
            "如果是扫描版 PDF，后续需要使用 OCR 识别。"
        )

    return text
